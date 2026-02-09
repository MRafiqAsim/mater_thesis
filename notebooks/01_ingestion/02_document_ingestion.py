# Databricks notebook source
# MAGIC %md
# MAGIC # 02 - Document Ingestion Pipeline (PDF, DOCX, XLSX, PPTX)
# MAGIC
# MAGIC **Phase 1: Setup & Ingestion | Week 1-2**
# MAGIC
# MAGIC This notebook handles the ingestion of various document formats into the Bronze zone.
# MAGIC
# MAGIC ## Supported Formats
# MAGIC - PDF (.pdf)
# MAGIC - Word (.docx, .doc)
# MAGIC - Excel (.xlsx, .xls, .xlsm)
# MAGIC - PowerPoint (.pptx, .ppt)
# MAGIC - Images with OCR (.jpg, .png) - via Azure Document Intelligence
# MAGIC
# MAGIC ## Data Flow
# MAGIC ```
# MAGIC RAW (files) → BRONZE (Delta: documents with extracted text)
# MAGIC ```

# COMMAND ----------

# MAGIC %md
# MAGIC ## 1. Install Dependencies

# COMMAND ----------

# MAGIC %pip install langchain langchain-community pypdf python-docx openpyxl python-pptx unstructured pandas pyarrow

# COMMAND ----------

dbutils.library.restartPython()

# COMMAND ----------

# MAGIC %md
# MAGIC ## 2. Configuration

# COMMAND ----------

# Get configuration from widgets or secrets
try:
    STORAGE_ACCOUNT = dbutils.widgets.get("storage_account")
    CONTAINER = dbutils.widgets.get("container")
except Exception:
    STORAGE_ACCOUNT = dbutils.secrets.get("azure-storage", "account-name")
    CONTAINER = dbutils.secrets.get("azure-storage", "container-name")

BASE_PATH = f"abfss://{CONTAINER}@{STORAGE_ACCOUNT}.dfs.core.windows.net"
RAW_PATH = f"{BASE_PATH}/raw"
BRONZE_PATH = f"{BASE_PATH}/bronze"

print(f"RAW_PATH: {RAW_PATH}")
print(f"BRONZE_PATH: {BRONZE_PATH}")

# COMMAND ----------

# MAGIC %md
# MAGIC ## 3. Define Document Schema

# COMMAND ----------

from pyspark.sql.types import StructType, StructField, StringType, IntegerType, TimestampType, ArrayType, LongType

DOCUMENT_SCHEMA = StructType([
    StructField("document_id", StringType(), False),
    StructField("filename", StringType(), True),
    StructField("file_path", StringType(), True),
    StructField("file_extension", StringType(), True),
    StructField("file_size_bytes", LongType(), True),
    StructField("page_count", IntegerType(), True),
    StructField("content", StringType(), True),
    StructField("content_length", IntegerType(), True),
    StructField("document_type", StringType(), True),  # pdf, word, excel, powerpoint
    StructField("title", StringType(), True),
    StructField("author", StringType(), True),
    StructField("created_date", StringType(), True),
    StructField("modified_date", StringType(), True),
    StructField("document_hash", StringType(), True),
    StructField("source_path", StringType(), True),
    StructField("ingestion_timestamp", TimestampType(), True),
])

# COMMAND ----------

# MAGIC %md
# MAGIC ## 4. Document Loaders

# COMMAND ----------

from typing import Dict, Any, List, Optional
from datetime import datetime
import hashlib
import os

def compute_file_hash(file_path: str) -> str:
    """Compute MD5 hash of file."""
    with open(file_path, 'rb') as f:
        return hashlib.md5(f.read()).hexdigest()

def load_pdf(file_path: str) -> Dict[str, Any]:
    """Load PDF document."""
    from langchain_community.document_loaders import PyPDFLoader

    loader = PyPDFLoader(file_path)
    pages = loader.load()

    content = "\n\n".join([p.page_content for p in pages])

    # Try to get metadata
    metadata = pages[0].metadata if pages else {}

    return {
        "content": content,
        "page_count": len(pages),
        "title": metadata.get("title"),
        "author": metadata.get("author"),
        "created_date": metadata.get("creationDate"),
        "modified_date": metadata.get("modDate"),
    }

def load_docx(file_path: str) -> Dict[str, Any]:
    """Load Word document."""
    from docx import Document

    doc = Document(file_path)

    # Extract text from paragraphs
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    content = "\n\n".join(paragraphs)

    # Extract from tables
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join([cell.text for cell in row.cells])
            if row_text.strip():
                content += "\n" + row_text

    # Get core properties
    props = doc.core_properties

    return {
        "content": content,
        "page_count": 1,  # DOCX doesn't have pages until rendered
        "title": props.title,
        "author": props.author,
        "created_date": props.created.isoformat() if props.created else None,
        "modified_date": props.modified.isoformat() if props.modified else None,
    }

def load_xlsx(file_path: str) -> Dict[str, Any]:
    """Load Excel document."""
    import openpyxl

    wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)

    content_parts = []
    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        content_parts.append(f"=== Sheet: {sheet_name} ===")

        for row in sheet.iter_rows(values_only=True):
            row_values = [str(cell) if cell is not None else "" for cell in row]
            if any(row_values):
                content_parts.append(" | ".join(row_values))

    content = "\n".join(content_parts)

    # Get properties
    props = wb.properties

    return {
        "content": content,
        "page_count": len(wb.sheetnames),
        "title": props.title,
        "author": props.creator,
        "created_date": props.created.isoformat() if props.created else None,
        "modified_date": props.modified.isoformat() if props.modified else None,
    }

def load_pptx(file_path: str) -> Dict[str, Any]:
    """Load PowerPoint document."""
    from pptx import Presentation

    prs = Presentation(file_path)

    content_parts = []
    for i, slide in enumerate(prs.slides, 1):
        content_parts.append(f"=== Slide {i} ===")

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                content_parts.append(shape.text)

    content = "\n\n".join(content_parts)

    # Get properties
    props = prs.core_properties

    return {
        "content": content,
        "page_count": len(prs.slides),
        "title": props.title,
        "author": props.author,
        "created_date": props.created.isoformat() if props.created else None,
        "modified_date": props.modified.isoformat() if props.modified else None,
    }

# COMMAND ----------

# MAGIC %md
# MAGIC ## 5. Document Loader Factory

# COMMAND ----------

LOADERS = {
    "pdf": load_pdf,
    "docx": load_docx,
    "doc": load_docx,
    "xlsx": load_xlsx,
    "xls": load_xlsx,
    "xlsm": load_xlsx,
    "pptx": load_pptx,
    "ppt": load_pptx,
}

def process_document(file_path: str, adls_path: str) -> Optional[Dict[str, Any]]:
    """
    Process a single document and return structured data.

    Args:
        file_path: Local file path (after download)
        adls_path: Original ADLS path

    Returns:
        Dictionary matching DOCUMENT_SCHEMA or None if error
    """
    filename = os.path.basename(file_path)
    ext = filename.split(".")[-1].lower() if "." in filename else ""

    if ext not in LOADERS:
        print(f"Unsupported format: {ext}")
        return None

    try:
        # Get file info
        file_size = os.path.getsize(file_path)
        file_hash = compute_file_hash(file_path)

        # Load document
        loader_func = LOADERS[ext]
        result = loader_func(file_path)

        return {
            "document_id": file_hash,
            "filename": filename,
            "file_path": file_path,
            "file_extension": ext,
            "file_size_bytes": file_size,
            "page_count": result.get("page_count", 1),
            "content": result.get("content", ""),
            "content_length": len(result.get("content", "")),
            "document_type": ext if ext not in ["doc", "xls", "ppt"] else {"doc": "word", "xls": "excel", "ppt": "powerpoint"}.get(ext, ext),
            "title": result.get("title"),
            "author": result.get("author"),
            "created_date": result.get("created_date"),
            "modified_date": result.get("modified_date"),
            "document_hash": file_hash,
            "source_path": adls_path,
            "ingestion_timestamp": datetime.utcnow(),
        }

    except Exception as e:
        print(f"Error processing {filename}: {e}")
        return None

# COMMAND ----------

# MAGIC %md
# MAGIC ## 6. Discover Documents in RAW Zone

# COMMAND ----------

def list_files_recursive(base_path: str, extensions: List[str]) -> List[str]:
    """List all files with given extensions recursively."""
    all_files = []

    def scan_dir(path: str):
        try:
            items = dbutils.fs.ls(path)
            for item in items:
                if item.isDir():
                    scan_dir(item.path)
                else:
                    ext = item.name.split(".")[-1].lower() if "." in item.name else ""
                    if ext in extensions:
                        all_files.append(item.path)
        except Exception as e:
            print(f"Error scanning {path}: {e}")

    scan_dir(base_path)
    return all_files

# Discover documents
SUPPORTED_EXTENSIONS = ["pdf", "docx", "doc", "xlsx", "xls", "xlsm", "pptx", "ppt"]
document_files = list_files_recursive(RAW_PATH, SUPPORTED_EXTENSIONS)

print(f"Found {len(document_files)} documents")

# Count by type
from collections import Counter
ext_counts = Counter([f.split(".")[-1].lower() for f in document_files])
print("\nBy file type:")
for ext, count in ext_counts.most_common():
    print(f"  {ext}: {count}")

# COMMAND ----------

# MAGIC %md
# MAGIC ## 7. Process Documents in Batches

# COMMAND ----------

def process_documents_batch(file_paths: List[str], batch_size: int = 50) -> tuple:
    """
    Process documents in batches.

    Returns:
        Tuple of (documents, errors)
    """
    documents = []
    errors = []

    for i, adls_path in enumerate(file_paths):
        try:
            # Download to local
            filename = adls_path.split("/")[-1]
            local_path = f"/tmp/doc_{i}_{filename}"
            dbutils.fs.cp(adls_path, f"file:{local_path}")

            # Process
            doc_data = process_document(local_path, adls_path)
            if doc_data:
                documents.append(doc_data)

            # Cleanup
            if os.path.exists(local_path):
                os.remove(local_path)

        except Exception as e:
            errors.append({"file": adls_path, "error": str(e)})

        # Progress
        if (i + 1) % batch_size == 0:
            print(f"Processed {i + 1}/{len(file_paths)} documents...")

    print(f"\nCompleted: {len(documents)} documents, {len(errors)} errors")
    return documents, errors

# Process all documents
if document_files:
    documents, errors = process_documents_batch(document_files)
else:
    documents, errors = [], []
    print("No documents found. Upload files to RAW zone first.")

# COMMAND ----------

# MAGIC %md
# MAGIC ## 8. Write to Bronze Zone

# COMMAND ----------

from pyspark.sql.functions import current_timestamp
from delta.tables import DeltaTable

def write_documents_to_bronze(documents: List[Dict], bronze_path: str):
    """Write documents to Bronze Delta Lake."""
    if not documents:
        print("No documents to write")
        return

    # Create DataFrame
    df = spark.createDataFrame(documents, schema=DOCUMENT_SCHEMA)
    df = df.withColumn("_ingested_at", current_timestamp())

    # Write to Delta
    table_path = f"{bronze_path}/documents"

    try:
        delta_table = DeltaTable.forPath(spark, table_path)
        delta_table.alias("target").merge(
            df.alias("source"),
            "target.document_id = source.document_id"
        ).whenMatchedUpdateAll().whenNotMatchedInsertAll().execute()
        print(f"Merged {df.count()} documents into {table_path}")
    except Exception:
        df.write.format("delta").mode("overwrite").save(table_path)
        print(f"Created new table with {df.count()} documents at {table_path}")

# Write to Bronze
if documents:
    write_documents_to_bronze(documents, BRONZE_PATH)

# COMMAND ----------

# MAGIC %md
# MAGIC ## 9. Verify Bronze Data

# COMMAND ----------

# Read and verify
docs_df = spark.read.format("delta").load(f"{BRONZE_PATH}/documents")

print(f"Total documents in Bronze: {docs_df.count()}")

# By document type
print("\nDocuments by type:")
docs_df.groupBy("document_type").count().display()

# COMMAND ----------

# Sample documents
display(docs_df.select("document_id", "filename", "document_type", "page_count", "content_length", "author").limit(10))

# COMMAND ----------

# MAGIC %md
# MAGIC ## 10. Content Statistics

# COMMAND ----------

from pyspark.sql.functions import avg, sum, min, max

# Content statistics
stats = docs_df.agg(
    avg("content_length").alias("avg_content_length"),
    min("content_length").alias("min_content_length"),
    max("content_length").alias("max_content_length"),
    sum("content_length").alias("total_characters"),
    avg("page_count").alias("avg_pages"),
    sum("page_count").alias("total_pages"),
)
display(stats)

# COMMAND ----------

# Size distribution
from pyspark.sql.functions import when

docs_df.withColumn(
    "size_category",
    when(docs_df.content_length < 1000, "small (<1K)")
    .when(docs_df.content_length < 10000, "medium (1K-10K)")
    .when(docs_df.content_length < 100000, "large (10K-100K)")
    .otherwise("very large (>100K)")
).groupBy("size_category").count().display()

# COMMAND ----------

# MAGIC %md
# MAGIC ## 11. Summary

# COMMAND ----------

summary = f"""
╔══════════════════════════════════════════════════════════════════╗
║            DOCUMENT INGESTION COMPLETE                           ║
╠══════════════════════════════════════════════════════════════════╣
║  Total Documents Processed  : {docs_df.count():<33} ║
║  Total Pages                : {docs_df.agg(sum("page_count")).collect()[0][0]:<33} ║
║  Total Characters           : {docs_df.agg(sum("content_length")).collect()[0][0]:<33} ║
║  Bronze Table Location      : {BRONZE_PATH}/documents           ║
╠══════════════════════════════════════════════════════════════════╣
║  NEXT STEPS:                                                     ║
║  1. Run 03_language_detection.py to detect EN/NL                 ║
║  2. Run chunking pipeline                                        ║
║  3. Proceed to NLP processing                                    ║
╚══════════════════════════════════════════════════════════════════╝
"""
print(summary)

# COMMAND ----------

# MAGIC %md
# MAGIC ## Next Notebook: Language Detection
# MAGIC
# MAGIC The next step is to detect the language (English/Dutch) of each document
# MAGIC for proper NLP model selection in Phase 2.
