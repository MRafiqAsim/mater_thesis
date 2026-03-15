"""
Document Parser

Parses various document formats (PDF, DOCX, XLSX, PPTX, TXT)
and extracts text content with metadata.
"""

import hashlib
import logging
import subprocess
import tempfile
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
from typing import List, Optional, Dict, Any, Union
from enum import Enum

logger = logging.getLogger(__name__)


class DocumentType(Enum):
    """Supported document types"""
    PDF = "pdf"
    DOCX = "docx"
    DOC = "doc"
    XLSX = "xlsx"
    XLS = "xls"
    PPTX = "pptx"
    PPT = "ppt"
    TXT = "txt"
    CSV = "csv"
    RTF = "rtf"
    HTML = "html"
    UNKNOWN = "unknown"


@dataclass
class ParsedDocument:
    """Parsed document with extracted content"""

    # Identification
    doc_id: str
    source_path: str
    doc_type: DocumentType

    # Content
    text: str
    pages: List[str] = field(default_factory=list)  # Text per page/sheet/slide
    tables: List[Dict[str, Any]] = field(default_factory=list)

    # Metadata
    title: Optional[str] = None
    author: Optional[str] = None
    created_date: Optional[datetime] = None
    modified_date: Optional[datetime] = None
    page_count: int = 0

    # Processing metadata
    extraction_time: datetime = field(default_factory=datetime.now)
    language: Optional[str] = None
    char_count: int = 0
    word_count: int = 0

    # Error tracking
    parse_errors: List[str] = field(default_factory=list)
    is_partially_parsed: bool = False

    def __post_init__(self):
        self.char_count = len(self.text)
        self.word_count = len(self.text.split())

    def to_dict(self) -> Dict[str, Any]:
        """Convert to dictionary"""
        return {
            "doc_id": self.doc_id,
            "source_path": self.source_path,
            "doc_type": self.doc_type.value,
            "title": self.title,
            "author": self.author,
            "created_date": self.created_date.isoformat() if self.created_date else None,
            "modified_date": self.modified_date.isoformat() if self.modified_date else None,
            "page_count": self.page_count,
            "char_count": self.char_count,
            "word_count": self.word_count,
            "language": self.language,
            "extraction_time": self.extraction_time.isoformat(),
            "has_tables": len(self.tables) > 0,
            "parse_errors": self.parse_errors,
        }


class DocumentParser:
    """
    Parse various document formats and extract text.

    Supports:
    - PDF: Text and table extraction
    - DOCX/DOC: Full document parsing
    - XLSX/XLS: Spreadsheet content
    - PPTX/PPT: Slide content
    - TXT/CSV/RTF/HTML: Plain text
    """

    EXTENSION_MAP = {
        ".pdf": DocumentType.PDF,
        ".docx": DocumentType.DOCX,
        ".doc": DocumentType.DOC,
        ".xlsx": DocumentType.XLSX,
        ".xls": DocumentType.XLS,
        ".pptx": DocumentType.PPTX,
        ".ppt": DocumentType.PPT,
        ".txt": DocumentType.TXT,
        ".csv": DocumentType.CSV,
        ".rtf": DocumentType.RTF,
        ".html": DocumentType.HTML,
        ".htm": DocumentType.HTML,
    }

    def __init__(
        self,
        extract_tables: bool = True,
        ocr_enabled: bool = False,
        max_file_size_mb: int = 100
    ):
        """
        Initialize the parser.

        Args:
            extract_tables: Whether to extract tables from documents
            ocr_enabled: Whether to use OCR for scanned PDFs
            max_file_size_mb: Maximum file size to process
        """
        self.extract_tables = extract_tables
        self.ocr_enabled = ocr_enabled
        self.max_file_size = max_file_size_mb * 1024 * 1024

    # Formats that msoffcrypto can handle
    ENCRYPTABLE_EXTENSIONS = {".xls", ".xlsx", ".doc", ".docx", ".ppt", ".pptx"}

    # Legacy formats that LibreOffice can convert
    LEGACY_CONVERSION_MAP = {
        ".ppt": ".pptx",
        ".doc": ".docx",
        ".xls": ".xlsx",
    }

    # Magic byte signatures for file type detection
    MAGIC_SIGNATURES = [
        (b"%PDF", ".pdf"),
        (b"PK\x03\x04", ".docx"),      # ZIP-based (DOCX/XLSX/PPTX — refined below)
        (b"\xd0\xcf\x11\xe0", ".msg"),  # OLE2 Compound Document (MSG/DOC/XLS/PPT)
        (b"{\\rtf", ".rtf"),
        (b"<html", ".html"),
        (b"<!DOCTYPE html", ".html"),
        (b"\x89PNG", ".png"),
        (b"\xff\xd8\xff", ".jpg"),
        (b"GIF8", ".gif"),
    ]

    def _detect_extensionless(self, path: Path) -> Path:
        """Detect file type from magic bytes for files without extensions."""
        if path.suffix:
            return path  # Already has an extension

        try:
            with open(path, "rb") as f:
                header = f.read(32)
        except Exception:
            return path

        if not header:
            return path

        detected_ext = None
        for magic, ext in self.MAGIC_SIGNATURES:
            if header.startswith(magic):
                detected_ext = ext
                break

        if not detected_ext:
            logger.info(f"Skipping extensionless file (unknown type): {path.name} | {path}")
            return path

        # Refine ZIP-based formats by checking internal content
        if detected_ext == ".docx":
            try:
                import zipfile
                with zipfile.ZipFile(path) as zf:
                    names = zf.namelist()
                    if any("word/" in n for n in names):
                        detected_ext = ".docx"
                    elif any("xl/" in n for n in names):
                        detected_ext = ".xlsx"
                    elif any("ppt/" in n for n in names):
                        detected_ext = ".pptx"
            except Exception:
                pass

        # Refine OLE2: MSG vs DOC/XLS/PPT
        if detected_ext == ".msg":
            try:
                import olefile
                if olefile.isOleFile(str(path)):
                    ole = olefile.OleFileIO(str(path))
                    streams = ole.listdir()
                    flat = ["/".join(s) for s in streams]
                    ole.close()
                    if any("__substg" in f.lower() for f in flat):
                        detected_ext = ".msg"  # Outlook message
                    elif any("worddocument" in f.lower() for f in flat):
                        detected_ext = ".doc"
                    elif any("workbook" in f.lower() for f in flat):
                        detected_ext = ".xls"
                    elif any("powerpoint" in f.lower() for f in flat):
                        detected_ext = ".ppt"
            except ImportError:
                pass  # olefile not installed — keep .msg as default for OLE2
            except Exception:
                pass

        # Rename file with detected extension
        new_path = path.with_suffix(detected_ext)
        try:
            import shutil
            shutil.copy2(str(path), str(new_path))
            logger.info(f"Detected extensionless file as {detected_ext}: {path.name} → {new_path.name} | {path}")
            return new_path
        except Exception as e:
            logger.warning(f"Failed to rename extensionless file {path.name}: {e}")
            return path

    def _decrypt_if_needed(self, path: Path) -> Path:
        """Decrypt Office files that use default (empty) password encryption."""
        if path.suffix.lower() not in self.ENCRYPTABLE_EXTENSIONS:
            return path

        try:
            import msoffcrypto
        except ImportError:
            return path

        try:
            with open(path, "rb") as f:
                office_file = msoffcrypto.OfficeFile(f)
                if not office_file.is_encrypted():
                    return path

                # Try empty password (common default encryption)
                decrypted_path = path.with_suffix(f".decrypted{path.suffix}")
                with open(decrypted_path, "wb") as out:
                    office_file.load_key(password="")
                    office_file.decrypt(out)

                logger.info(f"Decrypted {path.name} (default encryption)")
                return decrypted_path

        except Exception as e:
            logger.warning(f"Cannot decrypt {path.name}: {e}")
            return path

    def _convert_legacy(self, path: Path) -> Path:
        """Convert legacy Office formats (.ppt, .doc, .xls) using LibreOffice headless."""
        ext = path.suffix.lower()
        if ext not in self.LEGACY_CONVERSION_MAP:
            return path

        try:
            with tempfile.TemporaryDirectory() as tmp_dir:
                result = subprocess.run(
                    ["libreoffice", "--headless", "--convert-to",
                     self.LEGACY_CONVERSION_MAP[ext].lstrip("."),
                     "--outdir", tmp_dir, str(path)],
                    capture_output=True, text=True, timeout=60,
                )

                if result.returncode != 0:
                    logger.warning(f"LibreOffice conversion failed for {path.name}: {result.stderr.strip()}")
                    return path

                # Find converted file
                new_ext = self.LEGACY_CONVERSION_MAP[ext]
                converted = list(Path(tmp_dir).glob(f"*{new_ext}"))
                if not converted:
                    return path

                # Move converted file next to original
                output_path = path.with_suffix(new_ext)
                converted[0].rename(output_path)
                logger.info(f"Converted {path.name} → {output_path.name}")
                return output_path

        except FileNotFoundError:
            logger.warning("LibreOffice not installed — skipping legacy conversion")
            return path
        except Exception as e:
            logger.warning(f"Legacy conversion failed for {path.name}: {e}")
            return path

    def parse(self, file_path: Union[str, Path]) -> ParsedDocument:
        """
        Parse a document file.

        Pre-processes the file if needed:
        1. Decrypt Office files with default (empty) password
        2. Convert legacy formats (.ppt, .doc, .xls) via LibreOffice

        Args:
            file_path: Path to the document

        Returns:
            ParsedDocument with extracted content
        """
        path = Path(file_path)
        original_path = path

        if not path.exists():
            raise FileNotFoundError(f"File not found: {path}")

        # Check file size
        file_size = path.stat().st_size
        if file_size > self.max_file_size:
            raise ValueError(f"File too large: {file_size / 1024 / 1024:.1f} MB")

        # Pre-process: detect type for extensionless files, decrypt, convert legacy
        path = self._detect_extensionless(path)

        # Parse embedded email attachments (.msg)
        if path.suffix.lower() == ".msg":
            return self._parse_msg(path, original_path)

        path = self._decrypt_if_needed(path)
        path = self._convert_legacy(path)

        # Determine document type (from potentially converted file)
        doc_type = self._detect_type(path)

        # Generate document ID (from original path for consistency)
        doc_id = self._generate_id(original_path)

        # Parse based on type
        parser_map = {
            DocumentType.PDF: self._parse_pdf,
            DocumentType.DOCX: self._parse_docx,
            DocumentType.DOC: self._parse_doc,
            DocumentType.XLSX: self._parse_xlsx,
            DocumentType.XLS: self._parse_xls,
            DocumentType.PPTX: self._parse_pptx,
            DocumentType.PPT: self._parse_ppt,
            DocumentType.TXT: self._parse_txt,
            DocumentType.CSV: self._parse_csv,
            DocumentType.RTF: self._parse_rtf,
            DocumentType.HTML: self._parse_html,
        }

        parser = parser_map.get(doc_type, self._parse_txt)

        try:
            result = parser(path)
            result.doc_id = doc_id
            result.source_path = str(original_path)
            result.doc_type = self._detect_type(original_path)
            return result

        except Exception as e:
            logger.error(f"Error parsing {original_path}: {e}")
            # Return partial result
            return ParsedDocument(
                doc_id=doc_id,
                source_path=str(original_path),
                doc_type=self._detect_type(original_path),
                text="",
                parse_errors=[str(e)],
                is_partially_parsed=True
            )

    def parse_bytes(
        self,
        content: bytes,
        filename: str,
        doc_type: Optional[DocumentType] = None
    ) -> ParsedDocument:
        """
        Parse document from bytes.

        Args:
            content: Document content as bytes
            filename: Original filename (for type detection)
            doc_type: Override document type

        Returns:
            ParsedDocument
        """
        import tempfile

        # Determine type from filename if not specified
        if doc_type is None:
            ext = Path(filename).suffix.lower()
            doc_type = self.EXTENSION_MAP.get(ext, DocumentType.UNKNOWN)

        # Write to temp file and parse
        with tempfile.NamedTemporaryFile(suffix=Path(filename).suffix, delete=False) as f:
            f.write(content)
            temp_path = f.name

        try:
            result = self.parse(temp_path)
            result.source_path = filename  # Use original filename
            return result
        finally:
            Path(temp_path).unlink(missing_ok=True)

    def _detect_type(self, path: Path) -> DocumentType:
        """Detect document type from extension"""
        ext = path.suffix.lower()
        return self.EXTENSION_MAP.get(ext, DocumentType.UNKNOWN)

    def _generate_id(self, path: Path) -> str:
        """Generate unique document ID"""
        stat = path.stat()
        unique_str = f"{path}:{stat.st_size}:{stat.st_mtime}"
        return hashlib.md5(unique_str.encode()).hexdigest()[:16]

    # =========================================================================
    # PDF Parser
    # =========================================================================

    def _classify_pdf_page_fitz(self, page) -> dict:
        """
        Classify a single PDF page using PyMuPDF structural analysis.

        Returns dict with:
            type: 'digital', 'scanned', or 'hybrid'
            text_coverage: fraction of page area covered by text
            image_coverage: fraction of page area covered by images
            has_fonts: whether real fonts are used (not just outlines)
            text: extracted text (from fitz, higher quality than pypdf)
        """
        page_area = page.rect.width * page.rect.height
        if page_area == 0:
            return {"type": "scanned", "text_coverage": 0, "image_coverage": 0,
                    "has_fonts": False, "text": ""}

        # Signal 1: Text blocks with positions
        text_dict = page.get_text("dict", flags=0)
        text_area = 0.0
        char_count = 0
        has_fonts = False

        for block in text_dict.get("blocks", []):
            if block.get("type") == 0:  # text block
                bbox = block.get("bbox", (0, 0, 0, 0))
                block_area = (bbox[2] - bbox[0]) * (bbox[3] - bbox[1])
                text_area += block_area
                for line in block.get("lines", []):
                    for span in line.get("spans", []):
                        char_count += len(span.get("text", ""))
                        if span.get("font", ""):
                            has_fonts = True

        text_coverage = text_area / page_area

        # Signal 2: Image coverage
        image_area = 0.0
        images = page.get_images(full=True)
        for img in images:
            try:
                xref = img[0]
                img_rect = page.get_image_rects(xref)
                for rect in img_rect:
                    image_area += rect.width * rect.height
            except Exception:
                pass

        image_coverage = image_area / page_area

        # Signal 3: Extract text (fitz gives better quality than pypdf)
        page_text = page.get_text("text") or ""

        # Signal 4: Detect garbled glyph references
        import re
        glyph_count = len(re.findall(r'/g\d+', page_text))
        is_garbled = len(page_text) > 0 and glyph_count > len(page_text) / 10

        # Classification: combine signals
        has_real_text = char_count > 20 and has_fonts and not is_garbled and text_coverage > 0.02

        if has_real_text and image_coverage < 0.5:
            page_type = "digital"
        elif has_real_text and image_coverage >= 0.5:
            page_type = "hybrid"
        elif image_coverage > 0.3 or not has_real_text:
            page_type = "scanned"
        else:
            page_type = "digital"

        return {
            "type": page_type,
            "text_coverage": round(text_coverage, 3),
            "image_coverage": round(image_coverage, 3),
            "has_fonts": has_fonts,
            "text": page_text if not is_garbled else "",
        }

    def _parse_pdf(self, path: Path) -> ParsedDocument:
        """Parse PDF with PyMuPDF structural analysis and classification."""
        try:
            import fitz  # PyMuPDF
        except ImportError:
            logger.warning("PyMuPDF not installed — falling back to pypdf (no PDF classification)")
            return self._parse_pdf_fallback(path)

        doc = fitz.open(str(path))
        pages = []
        page_analyses = []
        metadata = {}

        # Extract metadata
        meta = doc.metadata or {}
        metadata = {
            "title": meta.get("title"),
            "author": meta.get("author"),
            "created": meta.get("creationDate"),
            "modified": meta.get("modDate"),
        }

        # Analyze each page
        for i in range(len(doc)):
            page = doc[i]
            analysis = self._classify_pdf_page_fitz(page)
            page_analyses.append(analysis)
            pages.append(analysis["text"])

        doc.close()

        # Classify entire document
        digital_count = sum(1 for a in page_analyses if a["type"] == "digital")
        scanned_count = sum(1 for a in page_analyses if a["type"] == "scanned")
        total = len(page_analyses)

        if total == 0:
            pdf_type = "scanned"
        elif scanned_count == 0:
            pdf_type = "digital"
        elif digital_count == 0:
            pdf_type = "scanned"
        else:
            pdf_type = "hybrid"

        # Build text output based on classification
        parse_errors = []
        if pdf_type == "scanned":
            logger.info(f"Scanned PDF: {path.name} ({total} pages) — needs OCR | {path}")
            full_text = ""
            parse_errors.append("scanned_pdf_needs_ocr")
        elif pdf_type == "hybrid":
            scanned_indices = [i for i, a in enumerate(page_analyses) if a["type"] == "scanned"]
            full_text = "\n\n".join(a["text"] for a in page_analyses)
            logger.info(f"Hybrid PDF: {path.name} — {len(scanned_indices)}/{total} pages need OCR | {path}")
            parse_errors.append(f"hybrid_pdf_ocr_pages:{','.join(str(i) for i in scanned_indices)}")
        else:
            full_text = "\n\n".join(a["text"] for a in page_analyses)

        # Extract tables (only for fully digital PDFs)
        tables = []
        if self.extract_tables and pdf_type == "digital":
            tables = self._extract_pdf_tables(path)

        # Parse dates
        created_date = self._parse_pdf_date(metadata.get("created"))
        modified_date = self._parse_pdf_date(metadata.get("modified"))

        result = ParsedDocument(
            doc_id="",
            source_path="",
            doc_type=DocumentType.PDF,
            text=full_text,
            pages=pages,
            tables=tables,
            title=metadata.get("title"),
            author=metadata.get("author"),
            created_date=created_date,
            modified_date=modified_date,
            page_count=total,
            parse_errors=parse_errors,
        )

        if pdf_type in ("scanned", "hybrid"):
            result.is_partially_parsed = True

        return result

    def _parse_pdf_fallback(self, path: Path) -> ParsedDocument:
        """Fallback PDF parser using pypdf (no structural classification)."""
        try:
            import pypdf
        except ImportError:
            try:
                import PyPDF2 as pypdf
            except ImportError:
                raise ImportError("Install pypdf or PyPDF2: pip install pypdf")

        text_parts = []
        pages = []
        metadata = {}

        with open(path, "rb") as f:
            reader = pypdf.PdfReader(f)

            if reader.metadata:
                metadata = {
                    "title": reader.metadata.get("/Title"),
                    "author": reader.metadata.get("/Author"),
                    "created": reader.metadata.get("/CreationDate"),
                    "modified": reader.metadata.get("/ModDate"),
                }

            for i, page in enumerate(reader.pages):
                try:
                    page_text = page.extract_text() or ""
                    pages.append(page_text)
                    text_parts.append(page_text)
                except Exception as e:
                    logger.warning(f"Error extracting page {i}: {e}")
                    pages.append("")

        tables = []
        if self.extract_tables:
            tables = self._extract_pdf_tables(path)

        created_date = self._parse_pdf_date(metadata.get("created"))
        modified_date = self._parse_pdf_date(metadata.get("modified"))

        return ParsedDocument(
            doc_id="",
            source_path="",
            doc_type=DocumentType.PDF,
            text="\n\n".join(text_parts),
            pages=pages,
            tables=tables,
            title=metadata.get("title"),
            author=metadata.get("author"),
            created_date=created_date,
            modified_date=modified_date,
            page_count=len(pages),
        )

    def _extract_pdf_tables(self, path: Path) -> List[Dict]:
        """Extract tables from PDF using tabula or camelot"""
        tables = []

        try:
            import tabula
            dfs = tabula.read_pdf(str(path), pages='all', silent=True)

            for i, df in enumerate(dfs):
                tables.append({
                    "table_id": i,
                    "rows": len(df),
                    "columns": list(df.columns),
                    "data": df.to_dict('records')
                })

        except ImportError:
            logger.debug("tabula not available, skipping table extraction")
        except Exception as e:
            logger.warning(f"Table extraction failed: {e}")

        return tables

    def _parse_pdf_date(self, date_str: Optional[str]) -> Optional[datetime]:
        """Parse PDF date format"""
        if not date_str:
            return None

        try:
            # Remove 'D:' prefix
            if date_str.startswith("D:"):
                date_str = date_str[2:]

            if len(date_str) >= 14:
                return datetime.strptime(date_str[:14], "%Y%m%d%H%M%S")
            elif len(date_str) >= 8:
                return datetime.strptime(date_str[:8], "%Y%m%d")
        except (ValueError, TypeError):
            pass

        return None

    # =========================================================================
    # DOCX Parser
    # =========================================================================

    def _parse_docx(self, path: Path) -> ParsedDocument:
        """Parse DOCX document"""
        try:
            from docx import Document
        except ImportError:
            raise ImportError("Install python-docx: pip install python-docx")

        doc = Document(str(path))
        text_parts = []
        tables = []

        # Extract paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)

        # Extract tables
        if self.extract_tables:
            for i, table in enumerate(doc.tables):
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)

                if table_data:
                    tables.append({
                        "table_id": i,
                        "rows": len(table_data),
                        "data": table_data
                    })

                    # Also add table text to content
                    table_text = "\n".join(["\t".join(row) for row in table_data])
                    text_parts.append(f"\n[Table {i+1}]\n{table_text}")

        # Extract metadata
        props = doc.core_properties

        return ParsedDocument(
            doc_id="",
            source_path="",
            doc_type=DocumentType.DOCX,
            text="\n\n".join(text_parts),
            tables=tables,
            title=props.title,
            author=props.author,
            created_date=props.created,
            modified_date=props.modified,
            page_count=1,  # DOCX doesn't have reliable page count
        )

    def _parse_doc(self, path: Path) -> ParsedDocument:
        """Parse legacy DOC document.

        Tries in order:
        1. python-docx (works for DOCX-compatible DOC files)
        2. antiword (reliable for legacy binary DOC)
        3. LibreOffice headless conversion (.doc → .docx)
        """
        # 1. Try python-docx
        try:
            return self._parse_docx(path)
        except Exception:
            pass

        # 2. Try antiword
        try:
            import subprocess
            result = subprocess.run(
                ["antiword", str(path)],
                capture_output=True, text=True, timeout=30,
            )
            if result.returncode == 0 and result.stdout.strip():
                return ParsedDocument(
                    doc_id="",
                    source_path="",
                    doc_type=DocumentType.DOC,
                    text=result.stdout.strip(),
                )
        except FileNotFoundError:
            logger.debug("antiword not installed, trying LibreOffice")
        except Exception as e:
            logger.debug(f"antiword DOC parsing failed: {e}")

        # 3. Try LibreOffice conversion (.doc → .docx)
        converted = self._convert_legacy(path)
        if converted != path and converted.exists():
            try:
                return self._parse_docx(converted)
            except Exception as e:
                logger.debug(f"LibreOffice-converted DOCX parsing failed: {e}")

        raise RuntimeError(f"Cannot parse DOC file: {path.name}. Install antiword or LibreOffice.")

    # =========================================================================
    # Excel Parser
    # =========================================================================

    def _parse_xlsx(self, path: Path) -> ParsedDocument:
        """Parse XLSX spreadsheet"""
        try:
            from openpyxl import load_workbook
        except ImportError:
            raise ImportError("Install openpyxl: pip install openpyxl")

        wb = load_workbook(str(path), read_only=True, data_only=True)
        text_parts = []
        pages = []  # Each sheet as a "page"
        tables = []

        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]

            # Skip chart-only sheets (no cell data)
            if not hasattr(sheet, 'iter_rows'):
                continue

            sheet_text = [f"=== Sheet: {sheet_name} ==="]
            sheet_data = []

            for row in sheet.iter_rows():
                row_values = []
                for cell in row:
                    value = cell.value
                    if value is not None:
                        row_values.append(str(value))
                    else:
                        row_values.append("")

                if any(row_values):
                    sheet_text.append("\t".join(row_values))
                    sheet_data.append(row_values)

            page_text = "\n".join(sheet_text)
            pages.append(page_text)
            text_parts.append(page_text)

            if sheet_data:
                tables.append({
                    "sheet_name": sheet_name,
                    "rows": len(sheet_data),
                    "data": sheet_data
                })

        # Metadata
        props = wb.properties

        wb.close()

        return ParsedDocument(
            doc_id="",
            source_path="",
            doc_type=DocumentType.XLSX,
            text="\n\n".join(text_parts),
            pages=pages,
            tables=tables,
            title=props.title if props else None,
            author=props.creator if props else None,
            created_date=props.created if props else None,
            modified_date=props.modified if props else None,
            page_count=len(pages),
        )

    def _parse_xls(self, path: Path) -> ParsedDocument:
        """Parse legacy XLS spreadsheet"""
        try:
            import xlrd
        except ImportError:
            raise ImportError("Install xlrd: pip install xlrd")

        wb = xlrd.open_workbook(str(path))
        text_parts = []
        pages = []
        tables = []

        for sheet in wb.sheets():
            sheet_text = [f"=== Sheet: {sheet.name} ==="]
            sheet_data = []

            for row_idx in range(sheet.nrows):
                row_values = [str(sheet.cell_value(row_idx, col_idx))
                             for col_idx in range(sheet.ncols)]

                if any(row_values):
                    sheet_text.append("\t".join(row_values))
                    sheet_data.append(row_values)

            page_text = "\n".join(sheet_text)
            pages.append(page_text)
            text_parts.append(page_text)

            if sheet_data:
                tables.append({
                    "sheet_name": sheet.name,
                    "rows": len(sheet_data),
                    "data": sheet_data
                })

        return ParsedDocument(
            doc_id="",
            source_path="",
            doc_type=DocumentType.XLS,
            text="\n\n".join(text_parts),
            pages=pages,
            tables=tables,
            page_count=len(pages),
        )

    # =========================================================================
    # PowerPoint Parser
    # =========================================================================

    def _parse_pptx(self, path: Path) -> ParsedDocument:
        """Parse PPTX presentation"""
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("Install python-pptx: pip install python-pptx")

        prs = Presentation(str(path))
        text_parts = []
        pages = []  # Each slide as a "page"

        for i, slide in enumerate(prs.slides, 1):
            slide_text = [f"=== Slide {i} ==="]

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)

                # Extract table content
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = [cell.text for cell in row.cells]
                        slide_text.append("\t".join(row_text))

            page_text = "\n".join(slide_text)
            pages.append(page_text)
            text_parts.append(page_text)

        # Metadata
        props = prs.core_properties

        return ParsedDocument(
            doc_id="",
            source_path="",
            doc_type=DocumentType.PPTX,
            text="\n\n".join(text_parts),
            pages=pages,
            title=props.title,
            author=props.author,
            created_date=props.created,
            modified_date=props.modified,
            page_count=len(pages),
        )

    def _parse_ppt(self, path: Path) -> ParsedDocument:
        """Parse legacy PPT presentation.

        Tries in order:
        1. python-pptx (works for PPTX-compatible PPT files)
        2. LibreOffice headless conversion (.ppt → .pptx)
        """
        # 1. Try python-pptx
        try:
            return self._parse_pptx(path)
        except Exception:
            pass

        # 2. Try LibreOffice conversion (.ppt → .pptx)
        converted = self._convert_legacy(path)
        if converted != path and converted.exists():
            try:
                return self._parse_pptx(converted)
            except Exception as e:
                logger.debug(f"LibreOffice-converted PPTX parsing failed: {e}")

        raise RuntimeError(f"Cannot parse PPT file: {path.name}. Install LibreOffice. | {path}")

    # =========================================================================
    # Text-based Parsers
    # =========================================================================

    def _parse_txt(self, path: Path) -> ParsedDocument:
        """Parse plain text file"""
        # Try different encodings
        text = ""
        for encoding in ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']:
            try:
                text = path.read_text(encoding=encoding)
                break
            except UnicodeDecodeError:
                continue

        return ParsedDocument(
            doc_id="",
            source_path="",
            doc_type=DocumentType.TXT,
            text=text,
            page_count=1,
        )

    def _parse_csv(self, path: Path) -> ParsedDocument:
        """Parse CSV file"""
        import csv

        text_parts = []
        tables = []
        table_data = []

        with open(path, 'r', encoding='utf-8', errors='replace') as f:
            reader = csv.reader(f)
            for row in reader:
                text_parts.append("\t".join(row))
                table_data.append(row)

        if table_data:
            tables.append({
                "rows": len(table_data),
                "columns": table_data[0] if table_data else [],
                "data": table_data
            })

        return ParsedDocument(
            doc_id="",
            source_path="",
            doc_type=DocumentType.CSV,
            text="\n".join(text_parts),
            tables=tables,
            page_count=1,
        )

    def _parse_rtf(self, path: Path) -> ParsedDocument:
        """Parse RTF file"""
        try:
            from striprtf.striprtf import rtf_to_text
            rtf_content = path.read_text(errors='replace')
            text = rtf_to_text(rtf_content)
        except ImportError:
            # Fallback: basic cleaning
            rtf_content = path.read_text(errors='replace')
            import re
            text = re.sub(r'\\[a-z]+\d*\s?', '', rtf_content)
            text = re.sub(r'[{}]', '', text)

        return ParsedDocument(
            doc_id="",
            source_path="",
            doc_type=DocumentType.RTF,
            text=text,
            page_count=1,
        )

    def _parse_html(self, path: Path) -> ParsedDocument:
        """Parse HTML file"""
        try:
            from bs4 import BeautifulSoup
            html = path.read_text(errors='replace')
            soup = BeautifulSoup(html, 'html.parser')

            # Remove script and style
            for tag in soup(['script', 'style', 'head']):
                tag.decompose()

            text = soup.get_text(separator='\n')
            title = soup.title.string if soup.title else None

        except ImportError:
            # Fallback: regex cleaning
            import re
            html = path.read_text(errors='replace')
            text = re.sub(r'<[^>]+>', ' ', html)
            text = re.sub(r'\s+', ' ', text)
            title = None

        return ParsedDocument(
            doc_id="",
            source_path="",
            doc_type=DocumentType.HTML,
            text=text.strip(),
            title=title,
            page_count=1,
        )

    def _parse_msg(self, path: Path, original_path: Path = None) -> ParsedDocument:
        """Parse Outlook .msg embedded email attachment."""
        original_path = original_path or path
        doc_id = self._generate_id(original_path)

        try:
            import extract_msg

            msg = extract_msg.Message(str(path))

            # Build text from email fields
            parts = []
            if msg.subject:
                parts.append(f"Subject: {msg.subject}")
            if msg.sender:
                parts.append(f"From: {msg.sender}")
            if msg.to:
                parts.append(f"To: {msg.to}")
            if msg.date:
                parts.append(f"Date: {msg.date}")
            parts.append("")  # blank line separator

            # Body: prefer plain text, fall back to HTML
            body = msg.body or ""
            if not body.strip() and msg.htmlBody:
                try:
                    from bs4 import BeautifulSoup
                    body = BeautifulSoup(msg.htmlBody, "html.parser").get_text(separator="\n")
                except ImportError:
                    import re
                    body = re.sub(r'<[^>]+>', ' ', msg.htmlBody.decode("utf-8", errors="replace"))

            parts.append(body)
            text = "\n".join(parts).strip()

            msg.close()

            logger.info(f"Parsed embedded email: {original_path.name} ({len(text)} chars) | {path}")

            return ParsedDocument(
                doc_id=doc_id,
                source_path=str(original_path),
                doc_type=DocumentType.UNKNOWN,
                text=text,
                title=msg.subject if msg.subject else original_path.name,
                page_count=1,
                metadata={"source_format": "msg", "embedded_email": True},
            )

        except Exception as e:
            logger.warning(f"Failed to parse MSG file: {original_path.name}: {e} | {path}")
            return ParsedDocument(
                doc_id=doc_id,
                source_path=str(original_path),
                doc_type=DocumentType.UNKNOWN,
                text="",
                parse_errors=[f"msg_parse_error: {e}"],
            )


# Convenience function
def parse_document(file_path: str) -> ParsedDocument:
    """
    Parse a document file.

    Args:
        file_path: Path to the document

    Returns:
        ParsedDocument with extracted content
    """
    parser = DocumentParser()
    return parser.parse(file_path)
